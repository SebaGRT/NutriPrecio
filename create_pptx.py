#!/usr/bin/env python3
"""Generate Sprint I PowerPoint presentation for NutriPrecio."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# NutriPrecio brand color
TEAL = RGBColor(0x00, 0x96, 0x88)
DARK_TEAL = RGBColor(0x00, 0x77, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
WARNING_ORANGE = RGBColor(0xFF, 0x98, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

DOCS_DIR = "docs"
OUTPUT_PATH = os.path.join(DOCS_DIR, "Sprint-I-NutriPrecio.pptx")


def add_background_shape(slide, color):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def add_colored_bar(slide, left, top, width, height, color):
    """Add a colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_text(slide, left, top, width, height, bullets, font_size=16,
                    font_color=DARK_GRAY, bullet_color=TEAL, font_name="Calibri"):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.level = 0
        p.space_after = Pt(8)
    return txBox


def create_title_slide(prs, title, subtitle, team, po, date):
    """Slide 1: Portada"""
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Top teal bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(0.15), TEAL)

    # Large teal accent bar on left
    add_colored_bar(slide, 0, Inches(1.8), Inches(0.25), Inches(3.5), TEAL)

    # Title
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11.5), Inches(1.2),
                 title, font_size=54, font_color=TEAL, bold=True,
                 font_name="Calibri Light")

    # Subtitle
    add_text_box(slide, Inches(0.6), Inches(2.9), Inches(11.5), Inches(0.6),
                 subtitle, font_size=24, font_color=MEDIUM_GRAY,
                 font_name="Calibri Light")

    # Separator line
    add_colored_bar(slide, Inches(0.6), Inches(3.6), Inches(4), Inches(0.02), TEAL)

    # Team info
    team_text = f"Equipo:\n{team}"
    add_text_box(slide, Inches(0.6), Inches(4.0), Inches(11.5), Inches(1.2),
                 team_text, font_size=16, font_color=DARK_GRAY, font_name="Calibri")

    # PO and Date
    po_date_text = f"Product Owner: {po}\n{date}"
    add_text_box(slide, Inches(0.6), Inches(5.3), Inches(11.5), Inches(0.8),
                 po_date_text, font_size=16, font_color=MEDIUM_GRAY,
                 font_name="Calibri")

    # Bottom teal bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)

    return slide


def create_intro_slide(prs):
    """Slide 2: Introduccion al Sprint I"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 "Introduccion al Sprint I", font_size=36, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    # Content
    bullets = [
        "Objetivo general: Plataforma de comparacion de precios de alimentos saludables en Chile",
        "Meta Sprint I: Implementar registro de tienda, login/registro de usuarios, y dashboard de vendedor",
        "Alcance: 3 historias de usuario (MV-03, MV-52, MV-54)",
        "9 tareas planificadas, 47 horas estimadas",
        "Duracion: Sprint de 2 semanas"
    ]
    add_bullet_text(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5),
                    bullets, font_size=20, font_color=DARK_GRAY)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def create_metodologia_slide(prs):
    """Slide 3: Metodologia y Enfoque Tecnico"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 "Metodologia y Enfoque Tecnico", font_size=36, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    # Two-column layout
    left_bullets = [
        "Framework: Scrum (Sprint de 2 semanas)",
        "Herramientas:",
        "  • Django 6.0 + DRF (backend)",
        "  • Angular 19 (frontend)",
        "  • Git (control de versiones)"
    ]
    add_bullet_text(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5),
                    left_bullets, font_size=18, font_color=DARK_GRAY)

    right_bullets = [
        "Roles del equipo:",
        "  • Product Owner: Ignacio Herrera",
        "  • Scrum Master: Sebastian Herrera",
        "  • Dev Team: Vicente Sepulveda, Matias Ramirez,",
        "              Benjamin Buzeta, Fernando Sepulveda",
        "",
        "Distribucion: Backend, Frontend, API,",
        "Autenticacion, Dashboard"
    ]
    add_bullet_text(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(4.5),
                    right_bullets, font_size=18, font_color=DARK_GRAY)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def create_diagram_slide(prs, title, image_name, explanation_bullets, image_top=Inches(1.4)):
    """Generic diagram slide with image and bullets."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 title, font_size=32, font_color=WHITE, bold=True,
                 font_name="Calibri Light")

    # Image
    img_path = os.path.join(DOCS_DIR, image_name)
    if os.path.exists(img_path):
        # Calculate dimensions to fit nicely
        pic = slide.shapes.add_picture(img_path, Inches(0.5), image_top,
                                        width=Inches(7.5))
    else:
        add_text_box(slide, Inches(0.5), image_top, Inches(7.5), Inches(0.5),
                     f"[Imagen no encontrada: {image_name}]",
                     font_size=14, font_color=WARNING_ORANGE)

    # Explanation bullets on the right
    add_bullet_text(slide, Inches(8.3), image_top, Inches(4.5), Inches(5),
                    explanation_bullets, font_size=16, font_color=DARK_GRAY)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def create_casos_uso_slide(prs):
    """Slide 4: Diagrama de Casos de Uso"""
    bullets = [
        "3 actores principales:",
        "  • Comprador: busca y compara productos",
        "  • Vendedor: registra tienda y gestiona precios",
        "  • Sistema: autentica y notifica",
        "11 casos de uso identificados",
        "Relaciones include/extend para",
        "reutilizacion de funcionalidades"
    ]
    return create_diagram_slide(prs, "Diagrama de Casos de Uso",
                                "diagrama-casos-de-uso.png", bullets)


def create_clases_slide(prs):
    """Slide 5: Diagrama de Clases"""
    bullets = [
        "5 entidades principales:",
        "  • User, Store, Product, Category, Price",
        "3 gestores de autenticacion y servicios",
        "1 interfaz para comparacion de precios",
        "Relaciones clave:",
        "  • Tienda-Usuario (owner)",
        "  • Producto-Categoria",
        "  • Precio-Producto-Tienda"
    ]
    return create_diagram_slide(prs, "Diagrama de Clases",
                                "diagrama-clases.png", bullets)


def create_actividades_slide(prs):
    """Slide 6: Diagrama de Actividades"""
    bullets = [
        "3 flujos principales implementados:",
        "  • MV-03: Registro de Tienda",
        "  • MV-52: Login y Registro",
        "  • MV-54: Dashboard del Vendedor",
        "Swimlanes: Usuario, Frontend,",
        "Backend, Base de Datos",
        "Muestra el flujo completo de",
        "interaccion del sistema"
    ]
    return create_diagram_slide(prs, "Diagrama de Actividades",
                                "diagrama-actividades.png", bullets,
                                image_top=Inches(1.4))


def create_review_slide(prs):
    """Slide 7: Sprint Review"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 "Sprint Review", font_size=36, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    # Left: Entregadas
    add_colored_bar(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.08),
                    ACCENT_GREEN)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Funcionalidades Entregadas", font_size=20,
                 font_color=ACCENT_GREEN, bold=True)

    entregadas = [
        "MV-52: Login y Registro (COMPLETO)",
        "  • Registro con encriptacion",
        "  • Tokens + interceptor HTTP",
        "  • AuthGuard de rutas",
        "",
        "MV-03: Registro de Tienda (COMPLETO)",
        "  • Modelo Store con owner FK",
        "  • CRUD /api/stores/",
        "  • Formulario frontend integrado",
        "",
        "MV-54: Dashboard Vendedor (COMPLETO)",
        "  • Perfil con avatar e indicadores",
        "  • Validacion is_seller",
        "  • Acciones rapidas"
    ]
    add_bullet_text(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
                    entregadas, font_size=14, font_color=DARK_GRAY)

    # Right: NO Entregadas
    add_colored_bar(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.08),
                    WARNING_ORANGE)
    add_text_box(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.4),
                 "No Entregadas / Parciales", font_size=20,
                 font_color=WARNING_ORANGE, bold=True)

    no_entregadas = [
        "Dashboard sidebar de navegacion",
        "  • Layout actual: pagina unica",
        "  • Sin barra lateral como se muestra",
        "    en el diagrama de actividades",
        "",
        "Estado real del perfil de tienda",
        "  • Estadisticas muestran '--'",
        "  • Endpoints disponibles, sin",
        "    consumo frontend en esta iteracion",
        "",
        "Gestion de productos desde dashboard",
        "  • Boton 'Configurar Perfil' deshabilitado",
        "  • CRUD de productos: fuera de alcance"
    ]
    add_bullet_text(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
                    no_entregadas, font_size=14, font_color=DARK_GRAY)

    # Metrics bar at bottom
    add_colored_bar(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.02),
                    LIGHT_GRAY)
    metrics_text = "Metricas: 7/9 tareas (78%)  |  ~35/47 horas (74%)  |  Velocidad: 74% del scope"
    add_text_box(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
                 metrics_text, font_size=14, font_color=MEDIUM_GRAY,
                 align=PP_ALIGN.CENTER)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def create_retrospective_slide(prs):
    """Slide 8: Sprint Retrospective"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 "Sprint Retrospective", font_size=36, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    # Left: Que salio bien
    add_colored_bar(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.08),
                    ACCENT_GREEN)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Que salio bien", font_size=20, font_color=ACCENT_GREEN, bold=True)

    bien = [
        "Distribucion clara de tareas",
        "  • Responsabilidades bien definidas",
        "  • Backend, frontend, API, auth,",
        "    dashboard: cada uno con owner",
        "",
        "Stack tecnologico apropiado",
        "  • Django + DRF + Angular 19",
        "  • Separacion de concerns facilito",
        "    desarrollo paralelo",
        "",
        "Autenticacion correcta desde inicio",
        "  • Tokens configurados dia 1",
        "  • Interceptor + AuthGuard funcionando",
        "",
        "Estructura organizada del proyecto"
    ]
    add_bullet_text(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
                    bien, font_size=14, font_color=DARK_GRAY)

    # Right: Que mejorar + Acciones
    add_colored_bar(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.08),
                    WARNING_ORANGE)
    add_text_box(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Que mejorar y Acciones", font_size=20,
                 font_color=WARNING_ORANGE, bold=True)

    mejorar = [
        "Falta de seguimiento diario",
        "  • Sin standups formales",
        "  • Tareas estancadas sin detectar",
        "",
        "Campo owner no incluido desde inicio",
        "  • Requirio migracion adicional",
        "",
        "No se valido is_seller inicialmente",
        "  • Compradores accedian al dashboard",
        "",
        "Sin seed data hasta el final",
        "  • Imposibilito demos intermedias",
        "",
        "Acciones proximo Sprint:",
        "  • Daily standups de 15 min",
        "  • Seed data desde dia 1",
        "  • Revisar UML antes de codificar"
    ]
    add_bullet_text(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
                    mejorar, font_size=14, font_color=DARK_GRAY)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def create_demo_slide(prs):
    """Slide 9: Transicion a Demo"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background_shape(slide, WHITE)

    # Header bar
    add_colored_bar(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
                 "Demo del Incremento", font_size=36, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    # Center content
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.6),
                 "Flujo de demostracion:", font_size=24,
                 font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    steps = [
        "1. Registrar vendedor  →  /register",
        "2. Login  →  token almacenado, redireccion a home",
        "3. Navegar al Dashboard  →  /dashboard",
        "4. Click 'Registrar Tienda'  →  /dashboard/store-form",
        "5. Completar formulario  →  POST /api/stores/",
        "6. Volver al Dashboard  →  confirmacion visual"
    ]
    add_bullet_text(slide, Inches(3.5), Inches(2.3), Inches(6.5), Inches(3),
                    steps, font_size=20, font_color=DARK_GRAY)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    # Backup note
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
                 "Backup: Screenshots disponibles si la demo en vivo falla",
                 font_size=18, font_color=MEDIUM_GRAY, italic=True,
                 align=PP_ALIGN.CENTER)

    # Bottom bar
    add_colored_bar(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), TEAL)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Portada
    create_title_slide(
        prs,
        title="NutriPrecio — Sprint I",
        subtitle="Certamen N°2 — Taller de Metodos de Innovacion Agiles",
        team="Vicente Sepulveda, Matias Ramirez, Benjamin Buzeta, Fernando Sepulveda, Sebastian Herrera (Scrum Master)",
        po="Ignacio Herrera",
        date="Mayo 2026"
    )

    # Slide 2: Introduccion
    create_intro_slide(prs)

    # Slide 3: Metodologia
    create_metodologia_slide(prs)

    # Slide 4: Casos de Uso
    create_casos_uso_slide(prs)

    # Slide 5: Clases
    create_clases_slide(prs)

    # Slide 6: Actividades
    create_actividades_slide(prs)

    # Slide 7: Sprint Review
    create_review_slide(prs)

    # Slide 8: Retrospective
    create_retrospective_slide(prs)

    # Slide 9: Demo
    create_demo_slide(prs)

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
